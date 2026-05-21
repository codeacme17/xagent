from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from xagent.web.api.files import _pptx_fallback_html


def test_pptx_fallback_html_escapes_slide_content(tmp_path: Path):
    deck_path = tmp_path / 'evil"><img src=x onerror=alert(1)>.pptx'
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "</h2><script>window.__xagent_xss=1</script>"
    table = slide.shapes.add_table(
        1, 1, Inches(1), Inches(2), Inches(4), Inches(1)
    ).table
    table.cell(0, 0).text = "<img src=x onerror=alert(1)>"
    slide.notes_slide.notes_text_frame.text = "<svg onload=alert(1)>"
    prs.save(str(deck_path))

    response = _pptx_fallback_html(deck_path)
    html = response.body.decode()

    assert "<script>" not in html
    assert "<img src=x" not in html
    assert "<svg onload" not in html
    assert "&lt;/h2&gt;&lt;script&gt;window.__xagent_xss=1&lt;/script&gt;" in html
    assert "&lt;img src=x onerror=alert(1)&gt;" in html
    assert "&lt;svg onload=alert(1)&gt;" in html
    assert "evil&quot;&gt;&lt;img src=x onerror=alert(1)&gt;.pptx" in html
